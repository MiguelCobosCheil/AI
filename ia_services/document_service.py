import os
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation


# -----------------------------
# PDF
# -----------------------------
def leer_pdf(ruta):
    try:
        reader = PdfReader(ruta)
        texto = ""
        for page in reader.pages:
            texto += page.extract_text() or ""
        return texto[:3000]
    except:
        return ""


# -----------------------------
# WORD
# -----------------------------
def leer_docx(ruta):
    try:
        doc = Document(ruta)
        texto = "\n".join([p.text for p in doc.paragraphs])
        return texto[:3000]
    except:
        return ""


# -----------------------------
# EXCEL / CSV
# -----------------------------
def leer_excel(ruta):
    try:
        df = pd.read_excel(ruta)
        return df.head(10).to_string()
    except:
        try:
            df = pd.read_csv(ruta)
            return df.head(10).to_string()
        except:
            return ""


# -----------------------------
# POWERPOINT
# -----------------------------
def leer_pptx(ruta):
    try:
        prs = Presentation(ruta)
        texto = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + "\n"
        return texto[:3000]
    except:
        return ""


# -----------------------------
# TXT
# -----------------------------
def leer_txt(ruta):
    try:
        with open(ruta, "r", encoding="utf-8") as f:
            return f.read()[:3000]
    except:
        return ""


# -----------------------------
# MAIN
# -----------------------------
def procesar_documento(ruta):

    ruta_lower = ruta.lower()

    if ruta_lower.endswith(".pdf"):
        return leer_pdf(ruta)

    if ruta_lower.endswith(".docx"):
        return leer_docx(ruta)

    if ruta_lower.endswith(".xlsx") or ruta_lower.endswith(".csv"):
        return leer_excel(ruta)

    if ruta_lower.endswith(".pptx"):
        return leer_pptx(ruta)

    if ruta_lower.endswith(".txt"):
        return leer_txt(ruta)

    return ""