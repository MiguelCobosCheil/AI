from fastapi import FastAPI, Body, Request
from fastapi.responses import FileResponse, RedirectResponse
from fastapi.staticfiles import StaticFiles
import requests
import json
import re
import time
import os
# os.environ["HF_HUB_OFFLINE"] = "1"
import base64
import pandas as pd
from dotenv import load_dotenv
from pypdf import PdfReader
from datetime import datetime
from PIL import Image
from fastapi import UploadFile, File, Form
from diffusers import StableDiffusionImg2ImgPipeline
from diffusers import StableDiffusionXLImg2ImgPipeline
from diffusers import StableDiffusionPipeline, StableDiffusionXLPipeline
import torch
import matplotlib.pyplot as plt
from openpyxl import Workbook
from urllib.parse import urljoin, urlparse
from io import StringIO
from ia_services.llm_service import llamar_ia, llamar_ia_json
from ia_services.video_service import generar_video
from ia_rules.rules_service import (
    crear_regla,
    cargar_reglas,
    obtener_regla,
    generar_contexto_regla
)
from bs4 import BeautifulSoup
from fastapi import FastAPI
from fastapi.staticfiles import StaticFiles

app = FastAPI()
app.mount("/desktop_assets", StaticFiles(directory="desktop_assets"), name="desktop_assets")

PIXEL_PATTERNS = {

    # -------------------------
    # GOOGLE / GMP
    # -------------------------
    "google": [
        "googletagmanager.com",
        "gtag(",
        "google-analytics.com",
        "analytics.js",
        "ga(",
        "gclid"
    ],
    "floodlight": [
        "doubleclick.net",
        "fls.doubleclick.net"
    ],
    "google_ads": [
        "googlesyndication.com",
        "adsbygoogle"
    ],
    "youtube": [
        "youtube.com",
        "ytag"
    ],

    # -------------------------
    # SOCIAL
    # -------------------------
    "meta": [
        "connect.facebook.net",
        "fbq(",
        "facebook.com/tr"
    ],
    "tiktok": [
        "tiktok.com",
        "ttq("
    ],
    "linkedin": [
        "snap.licdn.com",
        "linkedin.com/px"
    ],
    "twitter": [
        "static.ads-twitter.com",
        "twttr"
    ],

    # -------------------------
    # AFILIACIÓN
    # -------------------------
    "awin": ["awin1.com"],
    "appsflyer": ["appsflyer"],
    "blueknow": ["blueknow"],
    "youreko": ["youreko"],

    # -------------------------
    # PROGRAMMATIC
    # -------------------------
    "ttd": ["adsrvr.org"],
    "taboola": ["taboola"],
    "outbrain": ["outbrain"],
    "teads": ["teads"],
    "tribalfusion": ["tribalfusion"],
    "sonata": ["sonata"],
    "smoot": ["smoot"],
    "apexnative": ["apexnative"],


    # -------------------------
    # RETAIL
    # -------------------------
    "amazon_ads": ["amazon-adsystem"],

    # -------------------------
    # ANALYTICS / UX
    # -------------------------
    "contentsquare": ["contentsquare"],
    "clicktale": ["clicktale"],
    "trustarc": ["trustarc"],
    "medallia": ["medallia"],

    # -------------------------
    # CRM / CHAT
    # -------------------------
    "sprinklr": ["sprinklr"],

    # -------------------------
    # OTROS
    # -------------------------
    "firework": ["firework"],
    "funnkey": ["funnkey"],
    "insider": ["insider"],
    "criteo": ["criteo"]
}
RESOLUTIONS = {
    "low": (320, 180),      # ⚡ pruebas rápidas
    "sd": (768, 432),       # estándar
    "hd": (1280, 720),      # HD
    "fullhd": (1920, 1080)  # pesado
}
def es_pixel_valido(url):

    url = url.lower()

    # ❌ excluir basura
    if any(ext in url for ext in [
        ".png", ".jpg", ".jpeg", ".gif", ".svg", ".webp",
        ".css", ".woff", ".woff2", ".ttf"
    ]):
        return False

    # ❌ excluir consent / UI
    if any(k in url for k in ["consent", "notice", "cookie"]):
        return False

    # ✅ endpoints típicos de tracking
    if any(k in url for k in [
        "collect", "analytics", "track", "event",
        "pixel", "conversion", "doubleclick",
        "ads", "g/collect", "tr?"
    ]):
        return True

    return False

def analizar_har(path):

    with open(path, "r", encoding="utf-8") as f:
        har = json.load(f)

    encontrados = {}

    for entry in har.get("log", {}).get("entries", []):

        request = entry.get("request", {})
        response = entry.get("response", {})
        url = request.get("url", "").lower()

        if not url:
            continue

        # -------------------------
        # 🔥 FILTRO INTELIGENTE
        # -------------------------

        resource_type = entry.get("_resourceType", "")
        mime = response.get("content", {}).get("mimeType", "")
        size = response.get("content", {}).get("size", 0)

        # ❌ excluir recursos estáticos claros
        if any(ext in url for ext in [
            ".png", ".jpg", ".jpeg", ".gif", ".svg",
            ".css", ".woff", ".woff2", ".ttf"
        ]):
            continue

        # ❌ excluir consentimiento
        if any(k in url for k in ["consent", "notice", "cookie"]):
            continue

        # ❌ excluir HTML
        if "text/html" in mime:
            continue

        # ❌ excluir archivos grandes (no son pixels)
        if size and size > 200000:
            continue

        # -------------------------
        # 🔥 DETECCIÓN POR PROVEEDOR
        # -------------------------

        for proveedor, patrones in PIXEL_PATTERNS.items():

            if any(p in url for p in patrones):

                if proveedor not in encontrados:
                    encontrados[proveedor] = []

                encontrados[proveedor].append({
                    "url": url,
                    "method": request.get("method"),
                    "status": response.get("status"),
                    "type": resource_type
                })

    # -------------------------
    # limpiar duplicados
    # -------------------------
    for k in encontrados:
        seen = set()
        unique = []

        for item in encontrados[k]:
            if item["url"] not in seen:
                seen.add(item["url"])
                unique.append(item)

        encontrados[k] = unique

    return encontrados

def extraer_info_web(html, opciones):

    soup = BeautifulSoup(html, "html.parser")

    data = {}

    # -------------------------
    # UTIL: limpiar texto
    # -------------------------
    def get_text_clean():
        return " ".join(soup.get_text().split())[:2000]




    # -------------------------
    # TITLE
    # -------------------------
    if opciones.get("title") or opciones.get("titular"):
        data["title"] = (
            soup.title.string.strip()
            if soup.title and soup.title.string
            else ""
        )

    # -------------------------
    # META
    # -------------------------
    if opciones.get("meta"):
        metas = []
        for m in soup.find_all("meta"):
            metas.append({
                "name": m.get("name"),
                "content": m.get("content")
            })
        data["meta"] = metas

    # -------------------------
    # KEYWORDS (IA)
    # -------------------------
    if opciones.get("keywords"):
        texto = get_text_clean()
        data["keywords"] = llamar_ia(f"Saca palabras clave:\n{texto}")

    # -------------------------
    # RESUMEN IA
    # -------------------------
    if opciones.get("resumen"):
        texto = get_text_clean()
        data["resumen"] = llamar_ia(f"Resume esta página:\n{texto}")

    # -------------------------
    # SCHEMA
    # -------------------------
    if opciones.get("schema"):
        schemas = []
        for s in soup.find_all("script", type="application/ld+json"):
            if s.string:
                schemas.append(s.string)
        data["schema"] = schemas

    # -------------------------
    # PIXELS
    # -------------------------
        # NO hay bloque de pixels aquí
    # -------------------------
    # IMÁGENES
    # -------------------------
    if opciones.get("imagenes"):
        imgs = []
        for img in soup.find_all("img"):
            src = img.get("src")
            if src:
                imgs.append(src)
        data["imagenes"] = imgs[:10]

    # -------------------------
    # VIDEOS
    # -------------------------
    if opciones.get("videos"):
        videos = []
        for v in soup.find_all("video"):
            src = v.get("src")
            if src:
                videos.append(src)
        data["videos"] = videos

    # -------------------------
    # SEO + SCORING
    # -------------------------
    from web_services.seo_service import analizar_seo

    if opciones.get("seo"):
        seo_data = analizar_seo(soup)
        data.update(seo_data)
        score = 100
        errores = []

        # TITLE
        title = soup.title.string.strip() if soup.title and soup.title.string else ""

        if not title:
            score -= 20
            errores.append("Falta title")

        if len(title) > 65:
            score -= 5
            errores.append("Title demasiado largo")

        # META DESCRIPTION
        meta_desc = soup.find("meta", attrs={"name": "description"})
        if not meta_desc:
            score -= 20
            errores.append("Falta meta description")

        # CANONICAL
        canonical = soup.find("link", rel="canonical")
        if not canonical:
            score -= 10
            errores.append("Falta canonical")

        # ROBOTS
        robots = soup.find("meta", attrs={"name": "robots"})
        if robots and "noindex" in (robots.get("content") or "").lower():
            score -= 50
            errores.append("Página en noindex")

        # OG TAGS
        og_tags = soup.find_all("meta", property=lambda x: x and "og:" in x)
        if not og_tags:
            score -= 5
            errores.append("Sin Open Graph")

        # H1 / H2
        if not soup.find("h1"):
            score -= 10
            errores.append("Sin H1")

        if len(soup.find_all("h2")) < 1:
            score -= 5
            errores.append("Sin H2")

        # CONTENIDO
        texto = soup.get_text()
        if len(texto.strip()) < 500:
            score -= 20
            errores.append("Contenido pobre")

        # IMÁGENES ALT
        imagenes = soup.find_all("img")
        sin_alt = [img for img in imagenes if not img.get("alt")]

        if sin_alt:
            score -= 10
            errores.append(f"{len(sin_alt)} imágenes sin ALT")

        # HREFLANG
        hreflang = soup.find_all("link", rel="alternate")
        if not any("hreflang" in str(x) for x in hreflang):
            score -= 5
            errores.append("Sin hreflang")

        # ENLACES INTERNOS
        enlaces = soup.find_all("a", href=True)
        internos = [a for a in enlaces if "/" in a["href"]]

        if len(internos) < 5:
            score -= 5
            errores.append("Pocos enlaces internos")

        # KEYWORDS básicas
        palabras = texto.lower().split()
        data["top_keywords"] = list(set(palabras))[:10]

        data["score"] = max(score, 0)
        data["errores_seo"] = errores

    return data
# -----------------------------
# 🔥 MODELO RÁPIDO (SD 1.5)
# -----------------------------
print("Cargando modelo rápido...")
pipe_fast = None

def get_pipe_fast():
    global pipe_fast
    if pipe_fast is None:
        pipe_fast = StableDiffusionPipeline.from_pretrained(...).to("cpu")
    return pipe_fast

print("Modelo rápido listo")


# -----------------------------
# 🔥 MODELO PRO (SDXL BASE)
# -----------------------------
print("Cargando modelo SDXL (PRO)...")

pipe_sdxl = StableDiffusionXLPipeline.from_pretrained(
    "stabilityai/stable-diffusion-xl-base-1.0",
    local_files_only=True,
    torch_dtype=torch.float32
).to("cpu")

pipe_sdxl.enable_attention_slicing()

print("SDXL listo")


# -----------------------------
# 🔥 MODELO JUGGERNAUT (ULTRA REAL)
# -----------------------------
print("Cargando modelo JUGGERNAUT 🔥...")

pipe_juggernaut = None

def get_juggernaut():
    global pipe_juggernaut

    if pipe_juggernaut is None:
        print("🔥 Cargando Juggernaut bajo demanda...")

        pipe_juggernaut = StableDiffusionXLPipeline.from_pretrained(
            "RunDiffusion/Juggernaut-XL",
            torch_dtype=torch.float32,
            local_files_only=True
        ).to("cpu")

        pipe_juggernaut.enable_attention_slicing()  # ✅ AQUÍ

    return pipe_juggernaut

print("Juggernaut listo")


# -----------------------------
# 🔥 IMG2IMG (SD 1.5)
# -----------------------------
print("Cargando modelo img2img...")

pipe_img2img = StableDiffusionImg2ImgPipeline.from_pretrained(
    "runwayml/stable-diffusion-v1-5",
    local_files_only=True
).to("cpu")

print("Modelo img2img listo")


# -----------------------------
# 🔥 SDXL IMG2IMG
# -----------------------------
print("Cargando SDXL img2img...")

pipe_sdxl_img2img = StableDiffusionXLImg2ImgPipeline.from_pretrained(
    "stabilityai/stable-diffusion-xl-base-1.0",
    local_files_only=True,
    torch_dtype=torch.float32
).to("cpu")

pipe_sdxl_img2img.enable_attention_slicing()

print("SDXL img2img listo")


# -----------------------------
# 🔥 SDXL REFINER
# -----------------------------
print("Cargando SDXL refiner...")

pipe_refiner = StableDiffusionXLImg2ImgPipeline.from_pretrained(
    "stabilityai/stable-diffusion-xl-refiner-1.0",
    local_files_only=True,
    torch_dtype=torch.float32
).to("cpu")

pipe_refiner.enable_attention_slicing()

print("Refiner listo")


# -----------------------------
# ✅ TODOS LOS MODELOS CARGADOS
# -----------------------------
print("🚀 TODOS LOS MODELOS LISTOS")


# 🔥 IA
from ia_services.image_service import analizar_imagen_llava

# 🔥 DOCS
from docx import Document
from pptx import Presentation

load_dotenv()
@app.post("/web/export")
def exportar(data: dict):

    resultados = data.get("resultados", [])

    wb = Workbook()

    # -------------------------
    # HOJA RESUMEN
    # -------------------------
    ws_resumen = wb.active
    ws_resumen.title = "Resumen"
    ws_resumen.append(["URL", "Title", "Score"])

    # -------------------------
    # HOJAS EXTRA
    # -------------------------
    ws_seo = wb.create_sheet("SEO")
    ws_seo.append(["URL", "Errores"])

    ws_kw = wb.create_sheet("Keywords")
    ws_kw.append(["URL", "Keywords"])

    ws_meta = wb.create_sheet("Meta")
    ws_meta.append(["URL", "Name", "Content"])

    ws_schema = wb.create_sheet("Schema")
    ws_schema.append(["URL", "Schema"])

    # -------------------------
    # PIXELS
    # -------------------------
    ws_pixels = wb.create_sheet("Pixels")
    ws_pixels.append([
        "URL",
        "Proveedor",
        "Pixel URL",
        "Method",
        "Status",
        "Tipo"
    ])

    vistos = set()

    def detectar_proveedor(u):

        u = u.lower()

        if "doubleclick" in u or "googleadservices" in u:
            return "Google Ads / Floodlight"
        if "googletagmanager" in u:
            return "GTM"
        if "facebook" in u or "meta" in u:
            return "Meta"
        if "tiktok" in u:
            return "TikTok"
        if "linkedin" in u:
            return "LinkedIn"
        if "twitter" in u:
            return "Twitter"
        if "criteo" in u:
            return "Criteo"
        if "taboola" in u:
            return "Taboola"
        if "outbrain" in u:
            return "Outbrain"
        if "teads" in u:
            return "Teads"
        if "amazon" in u:
            return "Amazon Ads"
        if "sprinklr" in u:
            return "Sprinklr"
        if "trustarc" in u:
            return "TrustArc"

        return "Otros"

    # -------------------------
    # LOOP PRINCIPAL
    # -------------------------
    for item in resultados:

        url = item.get("url")
        data_item = item.get("data", {})

        # RESUMEN
        ws_resumen.append([
            url,
            data_item.get("title", ""),
            data_item.get("score", "")
        ])

        # SEO
        if data_item.get("errores_seo"):
            for err in data_item["errores_seo"]:
                ws_seo.append([url, err])

        # KEYWORDS
        if data_item.get("keywords"):
            ws_kw.append([url, data_item["keywords"]])

        # META
        for m in data_item.get("meta", []):
            ws_meta.append([
                url,
                m.get("name", ""),
                m.get("content", "")
            ])

        # SCHEMA
        for s in data_item.get("schema", []):
            ws_schema.append([url, str(s)[:300]])

        # PIXELS
        pixels = data_item.get("pixels", {})

        for proveedor, lista in pixels.items():

            for px in lista:

                if isinstance(px, dict):
                    px_url = px.get("url", "")
                    method = px.get("method", "")
                    status = px.get("status", "")
                    tipo = px.get("type", "")
                else:
                    px_url = str(px)
                    method = ""
                    status = ""
                    tipo = ""

                if not px_url:
                    continue

                key = (url, px_url)
                if key in vistos:
                    continue
                vistos.add(key)

                proveedor_final = detectar_proveedor(px_url)

                ws_pixels.append([
                    url,
                    proveedor_final,
                    px_url[:200],
                    method,
                    status,
                    tipo
                ])

    # -------------------------
    # GUARDAR (FUERA DEL LOOP)
    # -------------------------
    filename = f"export_{int(time.time())}.xlsx"
    path = f"files/{filename}"

    if not os.path.exists("files"):
        os.makedirs("files")

    wb.save(path)

    return {"ruta": filename}

@app.get("/files/{filename}")
def get_file(filename: str):

    ruta = os.path.join(os.getcwd(), "files", filename)

    if os.path.exists(ruta):
        return FileResponse(ruta)

    return {"error": "Archivo no encontrado"}

app.mount("/static", StaticFiles(directory="static"), name="static")

EMAIL_USER = os.getenv("EMAIL_USER")
EMAIL_PASS = os.getenv("EMAIL_PASS")

CLIENT_ID = os.getenv("CLIENT_ID")
TENANT_ID = os.getenv("TENANT_ID")
CLIENT_SECRET = os.getenv("CLIENT_SECRET")
REDIRECT_URI = os.getenv("REDIRECT_URI")

ACCESS_TOKEN = None
HISTORIAL_FILE = "historial.json"

EXTENSIONES_IMAGEN = [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff"]
EXTENSIONES_DOC = [".pdf", ".docx", ".xlsx", ".csv", ".pptx", ".txt"]

# -----------------------------
# Interpreta documento para crear imagen
# -----------------------------
def adaptar_documento_a_prompt(texto):

    prompt = f"""
Extrae SOLO la información útil para generar una imagen.

Texto:
{texto[:2000]}

Devuelve SOLO una descripción visual clara.
Ejemplo:
"hombre con camiseta roja, barba, fondo urbano"
"""

    try:
        return llamar_ia(prompt)
    except:
        return ""
# -----------------------------
# 🔥 BÚSQUEDA SEMÁNTICA (ARREGLADA SIN ROMPER NADA)
# -----------------------------
def es_relevante(texto, query):

    if not texto or not query:
        return False

    texto_lower = texto.lower()
    query_lower = query.lower()

    palabras = query_lower.split()
    if any(p in texto_lower for p in palabras):
        return True

    palabras_clave = [
        "presupuesto", "reforma", "obra", "coste", "precio",
        "albañileria", "electricidad", "fontaneria"
    ]

    if any(p in texto_lower for p in palabras_clave):
        return True

    prompt = f"""
Texto:
{texto[:1000]}

Consulta:
{query}

Responde SOLO:
SI -> si está relacionado con reformas, presupuestos, obras o costes
NO -> si no lo está
"""

    try:
        respuesta = llamar_ia(prompt, temperature=0).lower()
        return respuesta.startswith("si")

    except:
        return False


# -----------------------------
# LOGIN
# -----------------------------
@app.get("/login")
def login():
    url = (
        f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/authorize"
        f"?client_id={CLIENT_ID}"
        f"&response_type=code"
        f"&redirect_uri={REDIRECT_URI}"
        f"&response_mode=query"
        f"&scope=offline_access Mail.Read Mail.Send Mail.ReadWrite User.Read"
    )
    return RedirectResponse(url)


# -----------------------------
# CALLBACK
# -----------------------------
@app.get("/callback")
def callback(request: Request):

    global ACCESS_TOKEN

    code = request.query_params.get("code")

    token_url = f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/token"

    data = {
        "client_id": CLIENT_ID,
        "scope": "offline_access Mail.Read Mail.Send Mail.ReadWrite User.Read",
        "code": code,
        "redirect_uri": REDIRECT_URI,
        "grant_type": "authorization_code",
        "client_secret": CLIENT_SECRET,
    }

    r = requests.post(token_url, data=data)
    token_json = r.json()

    if "access_token" not in token_json:
        return {"error": "No access token"}

    ACCESS_TOKEN = token_json["access_token"]

    return RedirectResponse("/emails")


# -----------------------------
# HISTORIAL
# -----------------------------
def cargar_historial():
    if not os.path.exists(HISTORIAL_FILE):
        return []
    with open(HISTORIAL_FILE, "r") as f:
        return json.load(f)


def guardar_historial(data):
    with open(HISTORIAL_FILE, "w") as f:
        json.dump(data, f, indent=2)


# -----------------------------
# ADJUNTOS
# -----------------------------
def guardar_adjunto(att, carpeta="temp"):

    if not os.path.exists(carpeta):
        os.makedirs(carpeta)

    nombre = att.get("name")
    contenido = att.get("contentBytes")

    if not contenido:
        return None

    ruta = os.path.join(carpeta, nombre)

    with open(ruta, "wb") as f:
        f.write(base64.b64decode(contenido))

    return ruta


# -----------------------------
# LECTORES DOCUMENTOS
# -----------------------------
def leer_pdf(ruta):
    try:
        reader = PdfReader(ruta)
        texto = ""
        for page in reader.pages:
            texto += page.extract_text() or ""
        return texto[:3000]
    except:
        return ""


def leer_docx(ruta):
    try:
        doc = Document(ruta)
        return "\n".join([p.text for p in doc.paragraphs])[:3000]
    except:
        return ""


def leer_excel(ruta):
    try:
        xls = pd.ExcelFile(ruta)
        texto = ""

        for sheet in xls.sheet_names:
            df = xls.parse(sheet)

            if not df.empty:
                texto += df.to_string() + "\n"

        return texto[:3000]

    except Exception as e:
        print("ERROR EXCEL:", ruta, e)
        return ""


def leer_pptx(ruta):
    try:
        prs = Presentation(ruta)
        texto = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + "\n"
        return texto[:3000]
    except:
        return ""


def leer_txt(ruta):
    try:
        with open(ruta, "r", encoding="utf-8") as f:
            return f.read()[:3000]
    except:
        return ""

def leer_csv_robusto(ruta, header):
    for enc in ["utf-8", "latin-1", "cp1252"]:
        try:
            return pd.read_csv(
                ruta,
                header=header,
                sep=None,
                engine="python",
                encoding=enc
            )
        except:
            continue
    raise Exception("No se pudo leer el CSV con ningún encoding")
def procesar_documento(ruta):

    ruta_lower = ruta.lower()

    if ruta_lower.endswith(".pdf"):
        return leer_pdf(ruta)

    if ruta_lower.endswith(".docx"):
        return leer_docx(ruta)

    if ruta_lower.endswith(".xlsx") or ruta_lower.endswith(".csv"):
        return leer_excel(ruta)

    if ruta_lower.endswith(".pptx"):
        return leer_pptx(ruta)

    if ruta_lower.endswith(".txt"):
        return leer_txt(ruta)

    return ""


# -----------------------------
# RESUMEN DOCUMENTO
# -----------------------------
def resumir_documento(texto):

    if not texto:
        return ""

    prompt = f"Resume en 1 frase:\n{texto}"

    try:
        return llamar_ia(prompt)
    except:
        return texto[:200]


# -----------------------------
# IA EMAIL
# -----------------------------
def analizar_email(email, regla_id=None):

    contexto_imagenes = "\n".join(
        [f"- {img.get('descripcion')}" for img in email.get("imagenes", [])]
    )

    contexto_docs = "\n".join(
        [f"- {doc.get('resumen')}" for doc in email.get("documentos", [])]
    )

    # 🔥 CONTENIDO BASE DEL EMAIL
    contenido = f"""
EMAIL:
{email.get("cuerpo")}

DOCUMENTOS:
{contexto_docs}

IMÁGENES:
{contexto_imagenes}
"""

    # 🔵 CASO 1: SIN REGLA
    if not regla_id:

        prompt = f"""
EMAIL:
{email.get("cuerpo")}

DOCUMENTOS:
{contexto_docs}

IMÁGENES:
{contexto_imagenes}

Responde SOLO en JSON válido, sin texto adicional:

{{
  "requiere_respuesta": true/false,
  "respuesta": "respuesta útil"
}}
"""

    # 🟢 CASO 2: CON REGLA
    else:

        regla = obtener_regla(regla_id)
        contexto = generar_contexto_regla(regla)

        prompt = f"""
{contexto}

TAREA:
Evalúa si el siguiente email cumple las reglas.

CONTENIDO:
{contenido}

Responde SOLO en JSON válido, sin texto adicional:

{{
  "ok": true/false,
  "motivo": "explicación clara"
}}
"""

    # 🔥 EJECUCIÓN IA ROBUSTA
    try:
        resultado = llamar_ia_json(prompt)

        # 🛑 Validación básica
        if not resultado:
            raise Exception("JSON vacío")

        if regla_id:
            if "ok" not in resultado:
                raise Exception("JSON inválido (falta 'ok')")
        else:
            if "requiere_respuesta" not in resultado:
                raise Exception("JSON inválido (falta 'requiere_respuesta')")

        return resultado

    except Exception as e:
        print("❌ ERROR analizar_email:", e)

        # 🔴 fallback distinto según modo
        if regla_id:
            return {
                "ok": False,
                "motivo": "Error evaluando email"
            }
        else:
            return {
                "requiere_respuesta": True,
                "respuesta": "Estoy revisándolo."
            }
# -----------------------------
# 📁 ANALIZAR CARPETA (🔥 SOLO ESTA FUNCIÓN MODIFICADA)
# -----------------------------
# -----------------------------
# 📁 ANALIZAR CARPETA (🔥 SOLO ESTA FUNCIÓN MODIFICADA)
# -----------------------------
# 🔥 DENTRO DE analizar_carpeta (SOLO ESTA PARTE CAMBIADA)

@app.post("/analizar_carpeta")
def analizar_carpeta(data: dict = Body(...)):
    regla_id = data.get("regla_id")
    regla = obtener_regla(regla_id)
    ruta_base = data.get("ruta")

    if not ruta_base or not os.path.exists(ruta_base):
        return {"error": "Ruta no válida"}

    include_subfolders = data.get("include_subfolders", False)
    print("SUBFOLDERS:", include_subfolders)

    analizar_peso = data.get("peso", False)
    analizar_dimensiones = data.get("dimensiones", False)

    buscar_texto = data.get("buscar_texto", "")
    buscar_imagen = data.get("buscar_imagen", "")
    usar_busqueda_inteligente = data.get("usar_busqueda_inteligente", False)

    buscar_exacto = (data.get("buscar_exacto") or "").lower()

    nombre_archivo = (data.get("nombre_archivo") or "").lower()

    # 🔥 IMPORTANTE
    tipos_archivo = data.get("tipos_archivo") or []
    tipos_archivo = [t.lower() for t in tipos_archivo]

    fecha_desde = data.get("fecha_desde")
    fecha_hasta = data.get("fecha_hasta")

    print("TIPOS_BACK:", tipos_archivo)

    resultados = []
    total_archivos = 0

    hay_filtros = any([
        nombre_archivo,
        tipos_archivo,
        fecha_desde,
        fecha_hasta,
        buscar_exacto,
        (usar_busqueda_inteligente and (buscar_texto or buscar_imagen))
    ])

    if not include_subfolders and not hay_filtros:
        total = len([
            f for f in os.listdir(ruta_base)
            if os.path.isfile(os.path.join(ruta_base, f))
        ])
        return {"modo": "conteo", "total_archivos": total}

    for root, dirs, files in os.walk(ruta_base):

        if not include_subfolders:
            dirs.clear()

        for file in files:

            ruta = os.path.join(root, file)
            total_archivos += 1

            ext = os.path.splitext(file)[1].lower()

            print("DEBUG FILE:", file)
            print("EXT:", ext)

            # ✅ FILTRO TIPOS (ROBUSTO)
            if tipos_archivo:
                if ext not in tipos_archivo:
                    continue

            # ---------------------

            if nombre_archivo:
                if nombre_archivo not in file.lower():
                    continue

            fecha_archivo = datetime.fromtimestamp(os.path.getmtime(ruta))

            if fecha_desde:
                if fecha_archivo < datetime.fromisoformat(fecha_desde):
                    continue

            if fecha_hasta:
                if fecha_archivo > datetime.fromisoformat(fecha_hasta):
                    continue

            # ✅ BÚSQUEDA EXACTA
            if buscar_exacto:
                texto_total = file.lower()

                # 📄 DOCUMENTOS
                if ext in EXTENSIONES_DOC:
                    try:
                        texto_doc = procesar_documento(ruta)

                        print("\n📄 DOC:", file)
                        print(texto_doc[:500])  # 🔥 muestra contenido real
                        print("------")

                        if texto_doc:
                            texto_total += " " + texto_doc.lower()
                    except:
                        pass

                # 🖼️ IMÁGENES (IA)
                if ext in EXTENSIONES_IMAGEN:
                    try:
                        descripcion = analizar_imagen_llava(ruta_imagen=ruta, contexto=file)
                        if descripcion:
                            print("\n🖼️ IMG:", file)
                            print(descripcion)
                            print("------")
                            texto_total += " " + descripcion.lower()
                    except:
                        pass
                import re
                texto_total = re.sub(r"\s+", " ", texto_total.lower())
                busqueda = re.sub(r"\s+", " ", buscar_exacto.lower())

                palabras = busqueda.split()

                if not all(p in texto_total for p in palabras):
                    continue
              
            item = {"nombre": file, "ruta": ruta}

            if analizar_peso:
                item["peso"] = os.path.getsize(ruta)

            if analizar_dimensiones and ext in EXTENSIONES_IMAGEN:
                try:
                    img = Image.open(ruta)
                    item["dimensiones"] = f"{img.width}x{img.height}"
                except:
                    pass

            # 🔥 IA
            if usar_busqueda_inteligente and (buscar_texto or buscar_imagen):

                if ext in EXTENSIONES_IMAGEN and buscar_imagen:
                    try:
                        descripcion = analizar_imagen_llava(ruta_imagen=ruta, contexto=file)
                        if not es_relevante(descripcion, buscar_imagen):
                            continue
                        item["descripcion"] = descripcion
                    except:
                        continue

                elif ext in EXTENSIONES_DOC and buscar_texto:
                    texto = procesar_documento(ruta)
                    if not es_relevante(texto, buscar_texto):
                        continue
                    item["resumen"] = resumir_documento(texto)

         

            # -----------------------------
            # 🔥 EVALUACIÓN CON REGLA
            # -----------------------------
            if regla:

                contenido = item["nombre"]

                if "resumen" in item:
                    contenido += "\n" + item["resumen"]

                if "descripcion" in item:
                    contenido += "\n" + item["descripcion"]

                prompt_eval = f"""
            {generar_contexto_regla(regla)}

            TAREA:
            Evalúa si este archivo cumple las reglas.

            CONTENIDO:
            {contenido}

            Responde en JSON:
            {{
            "ok": true/false,
            "motivo": "explicación clara"
            }}
            """

                try:
                    item["evaluacion"] = llamar_ia_json(prompt_eval)

                except:
                    item["evaluacion"] = {
                        "ok": False,
                        "motivo": "Error evaluando archivo"
                    }

            resultados.append(item)

    return {
        "modo": "resultado",
        "total_archivos": total_archivos,
        "total_resultados": len(resultados),
        "resultados": resultados
    }
    
# -----------------------------
# 📩 ACCIONES (EMAILS)
# -----------------------------
@app.get("/acciones")
def acciones(regla_id: str = None):

    if not ACCESS_TOKEN:
        return {"error": "NO_AUTH"}

    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}

    r = requests.get(
        "https://graph.microsoft.com/v1.0/me/mailFolders/inbox/messages?$top=3",
        headers=headers
    )

    if r.status_code != 200:
        return {"error": "Graph error", "detail": r.text}

    data = r.json()
    resultado = []
    historial = cargar_historial()

    for msg in data.get("value", []):

        message_id = msg.get("id")

        attachments = []
        texto_adjuntos = ""
        documentos = []
        imagenes = []

        r_attach = requests.get(
            f"https://graph.microsoft.com/v1.0/me/messages/{message_id}/attachments",
            headers=headers
        )

        if r_attach.status_code == 200:

            for att in r_attach.json().get("value", []):

                if att.get("@odata.type") == "#microsoft.graph.fileAttachment":

                    ruta = guardar_adjunto(att)
                    if not ruta:
                        continue

                    nombre = att.get("name")

                    if any(nombre.lower().endswith(ext) for ext in EXTENSIONES_DOC):

                        texto = procesar_documento(ruta)
                        resumen = resumir_documento(texto)

                        documentos.append({"nombre": nombre, "resumen": resumen})
                        texto_adjuntos += resumen + "\n"

                    if any(nombre.lower().endswith(ext) for ext in EXTENSIONES_IMAGEN):

                        descripcion = analizar_imagen_llava(
                            ruta_imagen=ruta,
                            contexto=msg.get("bodyPreview")
                        )

                        imagenes.append({"nombre": nombre, "descripcion": descripcion})

                    attachments.append({"nombre": nombre, "ruta": ruta})

        email_obj = {
            "id": message_id,
            "asunto": msg.get("subject"),
            "cuerpo": msg.get("bodyPreview"),
            "preview": (msg.get("bodyPreview") or "")[:200],
            "texto_adjuntos": texto_adjuntos,
            "attachments": attachments,
            "imagenes": imagenes,
            "documentos": documentos
        }

        analisis = analizar_email(email_obj, regla_id)

        # 🔥 DECIDIR ACCIÓN SEGÚN MODO
        if "ok" in analisis:
            # 🟢 MODO REGLAS
            accion = "ok" if analisis["ok"] else "revisar"
            respuesta = analisis.get("motivo", "")
        else:
            # 🔵 MODO NORMAL
            accion = "responder" if analisis.get("requiere_respuesta") else "archivar"
            respuesta = analisis.get("respuesta", "")

        registro = {
            "message_id": message_id,
            "asunto": email_obj["asunto"],
            "preview": email_obj["preview"],
            "imagenes": imagenes,
            "documentos": documentos,
            "accion": accion,
            "respuesta": respuesta,
            "attachments": attachments,
            "fecha": datetime.now().isoformat(),
            "regla_id": regla_id  # 🔥 NUEVO
        }

        historial.append(registro)
        resultado.append(registro)

    guardar_historial(historial)
    return resultado


# -----------------------------
# FRONTENDS
# -----------------------------
@app.get("/folder")
def folder():
    return FileResponse("static/folder.html")


@app.get("/emails")
def emails():
    return FileResponse("static/emails.html")


@app.get("/")
def home():
    return FileResponse("static/index.html")

@app.get("/create")
def create():
    return FileResponse("static/create.html")

@app.get("/output.docx")
def descargar_word():
    return FileResponse("output.docx", filename="documento.docx")

@app.get("/output.pdf")
def descargar_pdf():
    return FileResponse("output.pdf", filename="documento.pdf")

@app.get("/output.xlsx")
def descargar_excel():
    return FileResponse("output.xlsx", filename="presupuesto.xlsx")

@app.get("/output.pptx")
def descargar_ppt():
    return FileResponse("output.pptx", filename="presentacion.pptx")


# -----------------------------
# OTROS
# -----------------------------
@app.get("/historial")
def ver_historial():
    return cargar_historial()


@app.get("/stats")
def stats():
    return cargar_historial()


@app.get("/debug")
def debug():

    if not ACCESS_TOKEN:
        return {"error": "no token"}

    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}

    r = requests.get(
        "https://graph.microsoft.com/v1.0/me/messages?$top=5",
        headers=headers
    )

    return {"status_code": r.status_code, "response": r.text}


def enviar_email(data):
    print("🚫 ENVÍO DESACTIVADO")

# -----------------------------
# ✨ GENERAR CONTENIDO (IA)
# -----------------------------
@app.post("/generar")
def generar(data: dict = Body(...)):

    prompt = data.get("prompt")
    regla_id = data.get("regla_id")

    if not prompt:
        return {"error": "Prompt vacío"}

    # 🔥 obtener regla
    regla = obtener_regla(regla_id)
    contexto = generar_contexto_regla(regla)

    # 🔥 construir prompt final
    if contexto:
        prompt_final = f"""
{contexto}

USUARIO:
{prompt}
"""
    else:
        prompt_final = prompt

    try:
        resultado = llamar_ia(prompt_final)
        return {"resultado": resultado}

    except Exception as e:
        return {"error": str(e)}

# -----------------------------
# 📄 GENERAR WORD
# -----------------------------
from docx import Document

@app.post("/generar_word")
def generar_word(data: dict = Body(...)):

    prompt = data.get("prompt")

    if not prompt:
        return {"error": "Prompt vacío"}

    try:
        # 🔥 1. Llamar a IA
        texto = llamar_ia(prompt)

        # 🔥 2. Crear Word
        doc = Document()
        doc.add_heading("Documento generado por IA", 0)
        doc.add_paragraph(texto)

        ruta = "output.docx"
        doc.save(ruta)

        return {"ruta": os.path.basename(ruta)}

    except Exception as e:
        return {"error": str(e)}      
    
# -----------------------------
# 📕 GENERAR PDF
# -----------------------------
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet

@app.post("/generar_pdf")
def generar_pdf(data: dict = Body(...)):

    prompt = data.get("prompt")

    if not prompt:
        return {"error": "Prompt vacío"}

    try:
        # IA
        texto = llamar_ia(prompt)

        # PDF
        doc = SimpleDocTemplate("output.pdf")
        styles = getSampleStyleSheet()

        content = []
        for linea in texto.split("\n"):
            content.append(Paragraph(linea, styles["Normal"]))

        doc.build(content)

        return {"ruta": "output.pdf"}

    except Exception as e:
        return {"error": str(e)}    
    
# -----------------------------
# 📊 GENERAR EXCEL
# -----------------------------
@app.post("/generar_excel")
def generar_excel(data: dict = Body(...)):

    prompt = data.get("prompt")

    if not prompt:
        return {"error": "Prompt vacío"}

    try:
        # IA (le pedimos estructura)
        texto = llamar_ia(f"""
Devuelve una tabla de presupuesto en formato:
Concepto | Precio

{prompt}
""")

        filas = []
        for linea in texto.split("\n"):
            if "|" in linea:
                partes = linea.split("|")
                if len(partes) >= 2:
                    filas.append({
                        "Concepto": partes[0].strip(),
                        "Precio": partes[1].strip()
                    })

        df = pd.DataFrame(filas)
        df.to_excel("output.xlsx", index=False)

        return {"ruta": "output.xlsx"}

    except Exception as e:
        return {"error": str(e)}    

from pptx import Presentation

@app.post("/generar_ppt")
def generar_ppt(data: dict = Body(...)):

    prompt = data.get("prompt")

    if not prompt:
        return {"error": "Prompt vacío"}

    try:
        texto = llamar_ia(prompt)

        prs = Presentation()

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Contenido generado"
        slide.placeholders[1].text = texto[:1000]

        prs.save("output.pptx")

        return {"ruta": "output.pptx"}

    except Exception as e:
        return {"error": str(e)}    
    


# -----------------------------
#  GENERAR IMAGENES
# -----------------------------
@app.post("/generar_imagen")
def generar_imagen(data: dict = Body(...)):

    prompt = data.get("prompt")
    modo = data.get("modo", "rapido")

    if not prompt:
        return {"error": "Prompt vacío"}

    try:

        if modo == "rapido":
            print("USANDO MODELO RAPIDO")
            image = pipe_fast(prompt, num_inference_steps=20).images[0]
        elif modo == "juggernaut":
            print("USANDO JUGGERNAUT 🔥")
            pipe = get_juggernaut()
            image = pipe_juggernaut(
                prompt,
                negative_prompt="blurry, low quality, distorted, bad anatomy, extra fingers, deformed face",
                num_inference_steps=100,
                guidance_scale=12
            ).images[0]
        elif modo == "ultra":
            print("USANDO SDXL ULTRA 🚀")

            image = pipe_sdxl(
                prompt,
                negative_prompt="blurry, low quality, distorted, bad anatomy",
                num_inference_steps=80,
                guidance_scale=9
            ).images[0]

        else:
            print("USANDO SDXL (PRO)")
            image = pipe_sdxl(
                prompt,
                negative_prompt="blurry, low quality, distorted, bad anatomy",
                num_inference_steps=25
            ).images[0]

        ruta = "output_image.png"
        image.save(ruta)

        return {"ruta": os.path.basename(ruta)}

    except Exception as e:
        return {"error": str(e)}
# -----------------------------
# 🖼️ GENERAR IMAGEN AVANZADA (con imagen o documento)
# -----------------------------
@app.post("/generar_imagen_avanzada")
async def generar_imagen_avanzada(
    prompt: str = Form(...),
    modo: str = Form("rapido"),
    imagen: UploadFile = File(None),
    documento: UploadFile = File(None)
):

    try:

        prompt_final = prompt

        # -------------------------
        # 📄 DOCUMENTO
        # -------------------------
        if documento:

            ruta_doc = f"temp_{documento.filename}"

            with open(ruta_doc, "wb") as f:
                f.write(await documento.read())

            texto_doc = procesar_documento(ruta_doc)

            print("\n📄 DOC USADO:")
            print(texto_doc[:500])

            descripcion = adaptar_documento_a_prompt(texto_doc)

            print("🧠 CONTEXTO EXTRAÍDO:", descripcion)

            prompt_final += f"\n\n{descripcion}"

        # -------------------------
        # 🖼️ IMAGEN
        # -------------------------
        if imagen:

            ruta_img = f"temp_{imagen.filename}"

            with open(ruta_img, "wb") as f:
                f.write(await imagen.read())

            img = Image.open(ruta_img).convert("RGB").resize((1024, 1024))

            print("\n🖼️ IMAGEN USADA:", ruta_img)
            print("MODO AVANZADO:", modo)

            if modo == "ultra":
                print("INPAINT ULTRA 🔥")

                import numpy as np

                image = pipe_sdxl_img2img(
                    prompt=prompt_final,
                    image=img,
                    strength=0.10,
                    num_inference_steps=160,
                    guidance_scale=15
                ).images[0]

                image = pipe_refiner(
                    prompt=prompt_final,
                    image=image,
                    num_inference_steps=50
                ).images[0]

            elif modo == "juggernaut":
                print("JUGGERNAUT CON CONTEXTO 🔥")
                pipe = get_juggernaut()
                image = pipe_juggernaut(
                    prompt_final,
                    negative_prompt="blurry, low quality, distorted, bad anatomy, extra fingers",
                    num_inference_steps=120,
                    guidance_scale=12
                ).images[0]

                image = pipe_refiner(
                    prompt=prompt_final,
                    image=image,
                    num_inference_steps=40
                ).images[0]            

            elif modo == "rapido":
                print("IMG2IMG RAPIDO")

                image = pipe_img2img(
                    prompt=prompt_final,
                    image=img,
                    strength=0.5,
                    num_inference_steps=20
                ).images[0]

            else:
                print("SDXL BASE GENERANDO...")

                image_base = pipe_sdxl_img2img(
                    prompt=prompt_final,
                    image=img,
                    strength=0.4,
                    num_inference_steps=30,
                    guidance_scale=7.5
                ).images[0]

                print("SDXL REFINER MEJORANDO...")

                image = pipe_refiner(
                    prompt=prompt_final,
                    image=image_base,
                    num_inference_steps=20
                ).images[0]

        else:
            # fallback sin imagen

            if modo == "ultra":
                print("ULTRA SIN IMAGEN (fallback limpio)")

                image = pipe_sdxl(
                    prompt_final,
                    num_inference_steps=80,
                    guidance_scale=9
                ).images[0]

            elif modo == "rapido":
                image = pipe_fast(prompt_final, num_inference_steps=20).images[0]

            else:
                image = pipe_sdxl(
                    prompt_final,
                    negative_prompt="blurry, low quality, distorted",
                    num_inference_steps=25
                ).images[0]
        
        ruta = "output_image_avanzada.png"
        image.save(ruta)

        return {"ruta": os.path.basename(ruta)}

    except Exception as e:
        return {"error": str(e)}  



@app.post("/generar_grafico")
async def generar_grafico(
    prompt: str = Form(""),
    tipo_grafico: str = Form("linea"),
    imagen: UploadFile = File(None),
    documento: UploadFile = File(None)
):

    try:
      

        df = None

        # -------------------------
        # 📄 DOCUMENTO
        # -------------------------
        if documento:

            ruta = f"temp_{documento.filename}"

            with open(ruta, "wb") as f:
                f.write(await documento.read())

            print("📄 DOCUMENTO:", ruta)

            # ✅ EXCEL / CSV
            if ruta.endswith(".xlsx") or ruta.endswith(".csv"):

                # 🔥 1. Leer archivo sin asumir nada
                if ruta.endswith(".xlsx"):
                    df_raw = pd.read_excel(ruta, header=None)
                else:
                   df_raw = leer_csv_robusto(ruta, header=None)

                # 🔥 2. Detectar fila header automáticamente
                header_row = None

                for i in range(min(10, len(df_raw))):
                    fila = df_raw.iloc[i].astype(str).str.lower()

                    if any(
                        "coste" in c or "precio" in c or "total" in c or "importe" in c or "€" in c
                        for c in fila
                    ):
                        header_row = i
                        break

                if header_row is None:
                    header_row = 0

                # 🔥 3. Volver a leer con header correcto
                if ruta.endswith(".xlsx"):
                    df = pd.read_excel(ruta, header=header_row)
                else:
                    df = leer_csv_robusto(ruta, header=header_row)

                print("HEADER DETECTADO:", header_row)
                print("COLUMNAS:", df.columns)
                print("📊 DATAFRAME INICIAL:")
                print(df.head(10))
                print("COLUMNAS:", df.columns)
                # 🔥 4. Detectar columna numérica automáticamente
                col_num = None
                max_validos = 0

                for col in df.columns:
                    serie = df[col].astype(str)

                    test = (
                        serie
                        .str.replace("€", "")
                        .str.replace(".", "")
                        .str.replace(",", ".")
                        .str.replace(r"[^\d.]", "", regex=True)
                    )

                    test_num = pd.to_numeric(test, errors="coerce")
                    validos = test_num.notna().sum()

                    print(f"COL {col} → validos: {validos}")

                    if validos > max_validos:
                        max_validos = validos
                        col_num = col
                        df[col] = test_num

                if col_num is None or max_validos == 0:
                    return {"error": "No se encontró columna numérica válida"}

                # 🔥 5. detectar columna de texto
                col_texto = None
                for col in df.columns:
                    if col != col_num:
                        col_texto = col
                        break

                if col_texto is None:
                    return {"error": "No se encontró columna de texto"}

                # 🔥 6. normalizar dataframe
                df = df[[col_texto, col_num]]
                df.columns = ["Concepto", "Valor"]

                df = df.dropna()
                df = df[df["Valor"] > 0]

                if df.empty:
                    return {"error": "No hay datos numéricos válidos"}
              
            # ✅ OTROS → IA
            else:

                texto = procesar_documento(ruta)

                print("🧠 TEXTO:", texto[:500])

                tabla = llamar_ia(f"""
Convierte este texto en datos económicos estructurados.

IMPORTANTE:
- Formato: Concepto | Valor
- El valor debe ser SOLO número (sin € ni texto)
- Ignora líneas sin números

Texto:
{texto[:2000]}
""")

                filas = []

                for linea in tabla.split("\n"):

                    linea = linea.strip()
                    if not linea:
                        continue

                    if "|" in linea:
                        partes = linea.split("|")
                    elif "," in linea:
                        partes = linea.split(",")
                    elif ":" in linea:
                        partes = linea.split(":")
                    else:
                        continue

                    if len(partes) >= 2:
                        concepto = partes[0].strip()
                        valor_raw = partes[1]

                        valor_limpio = re.sub(r"[^\d.]", "", valor_raw)

                        if valor_limpio:
                            filas.append({
                                "Concepto": concepto,
                                "Valor": valor_limpio
                            })

                df = pd.DataFrame(filas)

        # -------------------------
        # 🖼️ IMAGEN → IA
        # -------------------------
        elif imagen:

            ruta = f"temp_{imagen.filename}"

            with open(ruta, "wb") as f:
                f.write(await imagen.read())

            print("🖼️ IMAGEN:", ruta)

            descripcion = analizar_imagen_llava(ruta_imagen=ruta)

            tabla = llamar_ia(f"""
Convierte esto en datos numéricos.

Formato:
Concepto | Valor (solo número)

{descripcion}
""")

            filas = []

            for linea in tabla.split("\n"):

                linea = linea.strip()
                if not linea:
                    continue

                if "|" in linea:
                    partes = linea.split("|")
                elif "," in linea:
                    partes = linea.split(",")
                elif ":" in linea:
                    partes = linea.split(":")
                else:
                    continue

                if len(partes) >= 2:
                    concepto = partes[0].strip()
                    valor_raw = partes[1]

                    valor_limpio = re.sub(r"[^\d.]", "", valor_raw)

                    if valor_limpio:
                        filas.append({
                            "Concepto": concepto,
                            "Valor": valor_limpio
                        })

            df = pd.DataFrame(filas)

        # -------------------------
        # 📝 PROMPT
        # -------------------------
        elif prompt:

            tabla = llamar_ia(f"""
Genera datos numéricos.

Formato:
Concepto | Valor (solo número)

{prompt}
""")

            filas = []

            for linea in tabla.split("\n"):

                linea = linea.strip()
                if not linea:
                    continue

                if "|" in linea:
                    partes = linea.split("|")
                elif "," in linea:
                    partes = linea.split(",")
                elif ":" in linea:
                    partes = linea.split(":")
                else:
                    continue

                if len(partes) >= 2:
                    concepto = partes[0].strip()
                    valor_raw = partes[1]

                    valor_limpio = re.sub(r"[^\d.]", "", valor_raw)

                    if valor_limpio:
                        filas.append({
                            "Concepto": concepto,
                            "Valor": valor_limpio
                        })

            df = pd.DataFrame(filas)

        # -------------------------
        # VALIDACIÓN
        # -------------------------
        if df is None or df.empty:
            return {"error": "No hay datos válidos"}

        if len(df.columns) < 2:
            return {"error": "Formato incorrecto"}

        # limpieza numérica
        df[df.columns[1]] = pd.to_numeric(df[df.columns[1]], errors="coerce")
        df = df.dropna()

        # eliminar valores no válidos
        df = df[df[df.columns[1]] > 0]

        if df.empty:
            return {"error": "No hay datos numéricos válidos"}

        x = df[df.columns[0]].astype(str).tolist()
        y = df[df.columns[1]].astype(float).tolist()

        print("📊 DF FINAL LIMPIO:")
        print(df.head())

        # -------------------------
        # 📊 GRÁFICO
        # -------------------------
        import numpy as np

        plt.figure()

        indices = np.arange(len(x))

        if tipo_grafico == "barras":
            plt.bar(indices, y)
            plt.xticks(indices, x, rotation=45)

        elif tipo_grafico == "linea":
            plt.plot(indices, y, marker="o")
            plt.xticks(indices, x, rotation=45)

        elif tipo_grafico == "pie":
            plt.pie(y, labels=x, autopct="%1.1f%%")

        else:
            plt.plot(indices, y)
            plt.xticks(indices, x, rotation=45)

        plt.title("Distribución de costes")
        plt.ylabel("€")

        plt.tight_layout()

        ruta = f"grafico_{int(time.time())}.svg"

        plt.savefig(ruta, format="svg")
        plt.close()

        return {"ruta": os.path.basename(ruta)}

    except Exception as e:
        return {"error": str(e)}


# -------------------------
# CREAR REGLAS
# -------------------------
@app.post("/reglas/crear")
async def crear_regla_endpoint(
    nombre: str = Form(...),
    prompt_base: str = Form(""),
    documentos: list[UploadFile] = File(default=[]),
    imagenes: list[UploadFile] = File(default=[])
):

    textos_docs = []
    textos_imgs = []

    # -------------------------
    # 📄 DOCUMENTOS (MULTI)
    # -------------------------
    for doc in documentos:

        try:
            ruta = f"temp_{doc.filename}"

            with open(ruta, "wb") as f:
                f.write(await doc.read())

            texto = procesar_documento(ruta)

            if texto:
                textos_docs.append(texto[:2000])

        except Exception as e:
            print("ERROR DOC:", e)

    # -------------------------
    # 🖼️ IMÁGENES (MULTI)
    # -------------------------
    for img in imagenes:

        try:
            ruta = f"temp_{img.filename}"

            with open(ruta, "wb") as f:
                f.write(await img.read())

            descripcion = analizar_imagen_llava(ruta_imagen=ruta)

            if descripcion:
                textos_imgs.append(descripcion)

        except Exception as e:
            print("ERROR IMG:", e)

    # -------------------------
    # CREAR REGLA
    # -------------------------
    regla = crear_regla(
        nombre,
        prompt_base,
        textos_docs,
        textos_imgs
    )

    return {
        "ok": True,
        "regla": regla,
        "docs_procesados": len(textos_docs),
        "imagenes_procesadas": len(textos_imgs)
    } 

@app.get("/reglas")
def listar_reglas():
    return cargar_reglas()

@app.get("/dashboard")
def dashboard():
    return FileResponse("static/dashboard.html")

@app.get("/reglas_ui")
def reglas_ui():
    return FileResponse("static/reglas.html")

@app.post("/reglas/test")
async def test_regla(
    regla_id: str = Form(...),
    documento: UploadFile = File(None),
    texto: str = Form("")
):

    regla = obtener_regla(regla_id)

    if not regla:
        return {"error": "Regla no encontrada"}

    contenido = texto

    # 📄 documento opcional
    if documento:

        ruta = f"temp_{documento.filename}"

        with open(ruta, "wb") as f:
            f.write(await documento.read())

        contenido = procesar_documento(ruta)

    contexto = generar_contexto_regla(regla)

    prompt = f"""
{contexto}

TAREA:
Evalúa este contenido según la regla.

CONTENIDO:
{contenido}

Responde en JSON:
{{
  "ok": true/false,
  "motivo": "explicación clara"
}}
"""

    try:
        resultado = llamar_ia_json(prompt)

        return {
            "input": contenido[:1000],
            "resultado": resultado
        }

    except Exception as e:
        return {"error": str(e)}
    
@app.post("/reglas/feedback")
def feedback(data: dict = Body(...)):

    regla_id = data.get("regla_id")
    correcto = data.get("correcto")
    input_data = data.get("input")
    output = data.get("output")
    comentario = data.get("comentario", "")

    reglas = cargar_reglas()

    for r in reglas:
        if r["id"] == regla_id:

            if "feedback" not in r:
                r["feedback"] = []

            r["feedback"].append({
                "input": input_data,
                "output": output,
                "correcto": correcto,
                "comentario": comentario
            })

    # 🔥 IMPORTANTE → guardar
    with open("reglas.json", "w") as f:
        json.dump(reglas, f, indent=2)

    return {"ok": True}   

@app.delete("/reglas/{regla_id}")
def eliminar_regla(regla_id: str):

    reglas = cargar_reglas()

    reglas = [r for r in reglas if r["id"] != regla_id]

    with open("ia_rules/reglas.json", "w") as f:
        json.dump(reglas, f, indent=2)

    return {"ok": True} 




# -------------------------
# BOooOooooOOOOOoooottts
# -------------------------

@app.get("/web")
def web_ui():
    return FileResponse("static/web.html")
# -------------------------
# WEB CRAWL
# -------------------------
@app.post("/web/crawl")
def crawl(data: dict):

    from web_services.crawler_service import ejecutar_crawl

    base_url = data.get("url")
    opciones = data.get("opciones", {})
    max_pages = data.get("max_pages", 50)

    return ejecutar_crawl(base_url, max_pages, opciones)


# -------------------------
# WEB ANALYZE (🔥 MEJORADO)
# -------------------------
@app.post("/web/analyze")
def analyze(data: dict):

    from web_services.analyze_service import ejecutar_analyze

    urls = data.get("urls", [])
    opciones = data.get("opciones", {})

    return ejecutar_analyze(urls, opciones)
# -------------------------
# WEB COMPARE (OPTIMIZADO)
# -------------------------
@app.post("/web/compare")
def compare(data: dict):

    urls = data.get("urls", [])
    resultados = []

    with sync_playwright() as p:

        browser = p.chromium.launch(headless=True)

        for url in urls:

            try:
                page = browser.new_page()

                page.goto(url, timeout=20000)
                page.wait_for_load_state("networkidle")

                html = page.content()

                info = extraer_info_web(html, {
                    "title": True,
                    "meta": True,
                    "seo": True
                })

                resultados.append({
                    "url": url,
                    "score": info.get("score", 0),
                    "title": info.get("title", ""),
                    "errores": info.get("errores_seo", [])
                })

                page.close()

            except Exception as e:
                resultados.append({
                    "url": url,
                    "error": str(e)
                })

        browser.close()

    resultados = sorted(resultados, key=lambda x: x.get("score", 0), reverse=True)

    return resultados


# -------------------------
# WEB AGENT (OK)
# -------------------------
# -------------------------
# WEB AGENT (🔥 FIXED)
@app.post("/web/agent")
def web_agent(data: dict):

    from web_services.agent_service import ejecutar_agent

    try:
        steps = data.get("steps", [])
        url = data.get("url")

        if not steps:
            return {"error": "No hay pasos definidos"}

        resultados = ejecutar_agent(url, steps)

        return {
            "ok": True,
            "total_steps": len(steps),
            "resultados": resultados
        }

    except Exception as e:

        print("❌ ERROR EN AGENT:", e)

        return {
            "ok": False,
            "error": str(e)
        }

# -------------------------
# PIXELS (ENDPOINT DIRECTO)
# -------------------------
@app.post("/web/pixels")
def get_pixels(data: dict):

    from web_services.pixel_service import ejecutar_pixels

    url = data.get("url")

    if not url:
        return {"error": "URL requerida"}

    return ejecutar_pixels(url)

# -------------------------
# 💾 FLOWS (GLOBAL - FUNCIONA EN TODOS LOS NAVEGADORES)
# -------------------------

from web_services.agent_service import (
    guardar_flujo,
    cargar_flujos,
    eliminar_flujo
)

@app.post("/flows/save")
def save_flow(payload: dict):

    nombre = payload.get("name")
    flow = payload.get("flow")

    if not nombre or not flow:
        return {"error": "Datos incompletos"}

    guardar_flujo(nombre, flow)

    return {"ok": True}


@app.get("/flows/list")
def list_flows():

    return cargar_flujos()


@app.post("/flows/delete")
def delete_flow(payload: dict):

    nombre = payload.get("name")

    if not nombre:
        return {"error": "Nombre requerido"}

    eliminar_flujo(nombre)

    return {"ok": True}



# -------------------------
# grabar pantalla 1
# -------------------------
@app.post("/web/agent/ai")
def web_agent_ai(data: dict):

    from web_services.agent_service import ejecutar_agent_con_ia

    prompt = data.get("prompt")

    if not prompt:
        return {"error": "Prompt vacío"}

    return ejecutar_agent_con_ia(prompt)
@app.get("/web_ai")
def web_ai():
    return FileResponse("static/web_ai.html")
# -------------------------
# grabar pantalla 2
# -------------------------
@app.post("/desktop/ai")
def desktop_ai(data: dict):

    from ia_services.llm_service import llamar_ia

    flow = data.get("flow")
    prompt = data.get("prompt")

    respuesta = llamar_ia(f"""
Tienes este flujo:

{flow}

Usuario quiere:
{prompt}

Devuelve flujo modificado en JSON.
""")

    return {"result": respuesta}

recorder = None

@app.post("/desktop/record/start")
def start_record():

    global recorder

    import desktop_services.recorder_service as recorder_module

    recorder = recorder_module.start_recording()

    return {"ok": True}

@app.post("/desktop/record/stop")
def stop_record():

    global recorder

    import desktop_services.recorder_service as recorder_module

    flow = recorder_module.stop_recording(recorder)

    return {"flow": flow}

@app.post("/desktop/run")
def run_flow(data: dict):

    from desktop_services.executor_service import ejecutar_flow

    flow = data.get("flow")

    ejecutar_flow(flow)

    return {"ok": True}

@app.get("/desktop")
def desktop_ui():
    return FileResponse("static/desktop.html")
@app.post("/desktop/add_step")
def add_step():

    from desktop_services.recorder_service import capturar_desde_portapapeles

    path = capturar_desde_portapapeles()

    if not path:
        return {"error": "No hay imagen en portapapeles"}

    return {
        "type": "vision_step",
        "image": path,
        "prompt": ""
    }

FLOWS_FILE = "desktop_flows.json"

@app.post("/desktop/save_flow")
def save_flow(data: dict):

    flows = {}

    if os.path.exists(FLOWS_FILE):
        with open(FLOWS_FILE) as f:
            flows = json.load(f)

    flows[data["nombre"]] = data

    with open(FLOWS_FILE, "w") as f:
        json.dump(flows, f, indent=2)

    return {"ok": True}


@app.get("/desktop/flows")
def get_flows():
    if not os.path.exists(FLOWS_FILE):
        return {}
    return json.load(open(FLOWS_FILE))


@app.get("/desktop/get_flow/{nombre}")
def get_flow(nombre: str):
    flows = json.load(open(FLOWS_FILE))
    return flows.get(nombre, {})


@app.delete("/desktop/delete_flow/{nombre}")
def delete_flow(nombre: str):

    flows = json.load(open(FLOWS_FILE))

    if nombre in flows:
        del flows[nombre]

    with open(FLOWS_FILE, "w") as f:
        json.dump(flows, f, indent=2)

    return {"ok": True}

import desktop_services.recorder_service as recorder_module

@app.get("/desktop/record/steps")
def get_steps():
    print("📦 STEPS ACTUALES:", recorder_module.recording)
    return {"steps": recorder_module.recording}

# -------------------------
# crear video
# -------------------------
@app.post("/generar_video")
async def generar_video_endpoint(
    prompt: str = Form(""),
    imagen: UploadFile = File(None),
    frames: int = Form(16),
    fps: int = Form(8),
    quality: str = Form(...)
):

    try:
        
        if not prompt and not imagen:
            return {"error": "Falta prompt o imagen"}

        # 🔒 límites seguros
        frames = min(max(frames, 8), 40)
        fps = min(max(fps, 4), 12)
        # 🔥 validar quality SIEMPRE
        if quality not in RESOLUTIONS:
            quality = "sd"

        print(f"🎯 QUALITY RECIBIDA: {quality}")
        # 🔥 RESOLUCIÓN DESDE QUALITY
        w, h = RESOLUTIONS.get(quality, (768, 432))
        print(f"📐 RESOLUCIÓN FINAL: {w}x{h}")
        print(f"🎬 FRAMES: {frames} | FPS: {fps}")
        ruta_img = "temp_video.png"

        # -------------------------
        # 🖼️ INPUT
        # -------------------------
        if imagen:
            with open(ruta_img, "wb") as f:
                f.write(await imagen.read())

            image = Image.open(ruta_img).convert("RGB")

        else:
            image = pipe_sdxl(
                prompt,
                num_inference_steps=25
            ).images[0]

        # 🔥 IMPORTANTE: NO mejorar la imagen (como querías)
        image = image.resize((w, h))
        image.save(ruta_img)

        # -------------------------
        # 🎬 VIDEO (USANDO SERVICE)
        # -------------------------
        from ia_services.video_service import generar_video_pro

        filename = generar_video_pro(
            prompt=prompt,
            ruta_imagen=ruta_img,
            frames=frames,
            fps=fps,
            width=w,
            height=h
        )

        return {"ruta": filename}

    except Exception as e:
        print("❌ ERROR VIDEO:", e)
        return {"error": str(e)}